#Imports for making editable PPT
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
#import for accessing assets
from pathlib import Path
#imports for text handling
import re
import os
#Imports for timestamping file name
from datetime import datetime

# Directory paths for assets
script_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(script_dir)

CERT_DIR = os.path.join(project_root, "assets", "certifications")
ICON_DIR = os.path.join(project_root, "assets", "icons")
IMAGE_DIR = os.path.join(project_root, "assets", "images")  # For investment highlights images

# Available certification images
certifications = ["sedex.png", "iso 9001.png", "iso 14001.png", "fair trade.png", "energy star.png", "TGA - Australia.jpeg", "FDA.png", "B corporation.png", "MHRA - UK.png", "Salesforce Certified Financial Services Cloud Accredited Professional.jpeg","Salesforce Certified Administrator.png", " International Safety Award from the British Safety Council.png", "IDMA Quality Excellence Certified.jpeg"]

# Available icon images
icons = ["Cyber_Security_Hacking_Safety.png","Delivery_Service.png","Ecofriendly_Plant.png","Factory_Manufacturing_Industry.png","Global_Internet.png","Graph_Growth.png","Inventory_Storage.png","Restaurant_Food.png","song.png","money.png" ]


highlight_images = ["Industry.jpeg" , "pharma.jpeg", "tech.jpeg", "transport_logistics.jpeg"]

def apply_formatting_to_placeholder(placeholder, text):
    text_frame = placeholder.text_frame
    text_frame.clear()
    
    # Handle lists (arrays) - each item becomes a paragraph
    if isinstance(text, list):
        for idx, item_text in enumerate(text):
            # Use existing first paragraph, add new ones for subsequent items
            if idx == 0:
                p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
            else:
                p = text_frame.add_paragraph()
            
            # Split by ** markers
            parts = re.split(r'\*\*(.*?)\*\*', item_text)
            
            for i, part in enumerate(parts):
                if part:
                    run = p.add_run()
                    run.text = part
                    run.font.bold = (i % 2 == 1)
    
    # Handle strings - single paragraph
    else:
        # Reuse the first paragraph instead of adding new one
        p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        
        # Split by ** markers
        parts = re.split(r'\*\*(.*?)\*\*', str(text))
        
        for i, part in enumerate(parts):
            if part:
                run = p.add_run()
                run.text = part
                run.font.bold = (i % 2 == 1)

#Creates PPT
def create_ppt(slides_data):
    
    def add_slide_from_template(prs, layout_index, placeholder_text=None, placeholder_charts=None, placeholder_images=None):
        
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        if placeholder_text:
            for idx, text_val in placeholder_text.items():
                try:
            # Pass text directly - the function handles both lists and strings
                    apply_formatting_to_placeholder(slide.placeholders[idx], text_val)
                except KeyError:
                    print(f"Placeholder {idx} not found on slide layout {layout_index}")
        if placeholder_charts:
            for idx, (chart_type, chart_data) in placeholder_charts.items():
                if not chart_data: 
                    continue
                try:
                    placeholder = slide.placeholders[idx]
                except KeyError:
                    print(f"Placeholder {idx} not found for chart.")
                    continue

                # Capture dimensions before deleting placeholder
                x, y, cx, cy = placeholder.left, placeholder.top, placeholder.width, placeholder.height

                # Delete placeholder to prevent "Add Chart" text from appearing
                sp = placeholder.element
                sp.getparent().remove(sp)

                # Create the new chart
                graphic_frame = slide.shapes.add_chart(
                    chart_type, x, y, cx, cy, chart_data
                )
                chart = graphic_frame.chart

                # --- CHART FORMATTING ---
                # 1. Add Legend
                chart.has_legend = True
                chart.legend.include_in_layout = False 

                # 2. Add Data Labels
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels
                
                # Specific formatting based on chart type
                if chart_type == XL_CHART_TYPE.PIE:
                    chart.has_title = False
                    data_labels.show_category_name = True  # Shows "Revenue", "Cost", etc.
                    data_labels.show_value = True          # Shows "100", "50", etc.
                    data_labels.show_percentage = False    # Set to True if you want %
                    data_labels.position = XL_LABEL_POSITION.BEST_FIT
                
                elif chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
                    data_labels.show_value = True  # Usually just value is enough for bars

        if placeholder_images:
            for idx, image_path in placeholder_images.items():
            # 1. If no path provided, DELETE the placeholder
                if not image_path:
                    delete_placeholder(slide, idx)
                    continue

            # 2. If file doesn't exist, DELETE the placeholder
                if not os.path.exists(image_path):
                    print(f"Image not found: {image_path}")
                    delete_placeholder(slide, idx)
                    continue

            # 3. Insert the picture
                try:
                    placeholder = slide.placeholders[idx]
                    # DEBUG LINE:
                    # print(f"Placeholder {idx} type: {placeholder.placeholder_format.type}, Shape type: {placeholder.shape_type}")
                    placeholder.insert_picture(image_path)
                except Exception as e:
                    # If insertion fails, DELETE the placeholder
                    print(f"Placeholder {idx} does not support images: {e}")
                    delete_placeholder(slide, idx)

        return slide

    def delete_slide(prs, slide_index):
        slide_id_list = prs.slides._sldIdLst
        slides = list(slide_id_list)

        prs.part.drop_rel(slides[slide_index].rId)
        slide_id_list.remove(slides[slide_index])

    def delete_placeholder(slide, idx):
        try:
            placeholder = slide.placeholders[idx]
            sp = placeholder.element
            sp.getparent().remove(sp)
        except (KeyError, AttributeError):
            pass
    
    bar_data = CategoryChartData()
    pie_data = CategoryChartData()

    #  PARSE PIE CHART DATA 
    pie_data_obj = slides_data.get("pie_chart_data", {})
    # Default fallback if empty or missing
    if not pie_data_obj:
        pie_data_obj = {"title": "Default Pie", "categories": ["A", "B"], "values": [50, 50]}
        
    pie_title = pie_data_obj.get("title", "Default Pie")
    pie_cats = pie_data_obj.get("categories", [])
    pie_vals = pie_data_obj.get("values", [])

    # Ensure equal length (avoid index errors)
    min_len = min(len(pie_cats), len(pie_vals))
    pie_data.categories = pie_cats[:min_len]
    pie_data.add_series(pie_title, pie_vals[:min_len])
    
    # PARSE BAR CHART DATA 
    bar_data_obj = slides_data.get("bar_chart_data", {})
    # Default fallback
    if not bar_data_obj:
        bar_data_obj = {"title": "Default Bar", "categories": ["A", "B"], "values": [50, 50]}

    bar_title = bar_data_obj.get("title", "Default Bar")
    bar_cats = bar_data_obj.get("categories", [])
    bar_vals = bar_data_obj.get("values", [])

    min_len_bar = min(len(bar_cats), len(bar_vals))
    bar_data.categories = bar_cats[:min_len_bar]
    bar_data.add_series(bar_title, bar_vals[:min_len_bar])

    # LOAD PRESENTATION TEMPLATE

    layout_path = os.path.join(project_root, "Layout.pptx")
    if not os.path.exists(layout_path):
        print(f"Error: Layout file not found at {layout_path}")
        # fallback or let it crash with a better message, but for now just print

    prs = Presentation(layout_path)

    # CERTIFICATIONS PROCESSING
    
    cert_raw = slides_data.get("certifications", "")
    cert = cert_raw.split("||") if cert_raw else []
    
    # Pad with empty strings if fewer than 3
    while len(cert) < 3:
        cert.append("")
    
    # Clean up filenames and build full paths
    cert = [os.path.join(CERT_DIR, c.strip()) if c.strip() else "" for c in cert]

    # ICONS PROCESSING
    icon_raw = slides_data.get("icons", "")
    icon_list = icon_raw.split("||") if icon_raw else []
    
    # Pad with empty strings if fewer than 5
    while len(icon_list) < 5:
        icon_list.append("")
    
    # Clean up filenames and build full paths
    icon_list = [os.path.join(ICON_DIR, i.strip()) if i.strip() else "" for i in icon_list]
    #Debug
    # print("Icon paths:")
    # for idx, icon_path in enumerate(icon_list):
    #     if icon_path:
    #         exists = os.path.exists(icon_path)
    #         # print(f"  Icon {idx}: {icon_path} - Exists: {exists}")

    
    sector = slides_data.get("sector", "B2B") 
    
    num_layouts = len(prs.slide_layouts)
    # print(f"DEBUG: Found {num_layouts} layouts in template.")

    # Default for B2B, Manufacturing, Tech, Pharma, Logistics
    profile_layout_idx = 1
    financials_layout_idx = 2
        
    # SAFETY CHECK: If layout index is out of range, fallback to default
    if profile_layout_idx >= num_layouts:
        print(f"Warning: Layout {profile_layout_idx} not found. Fallback to 1.")
        profile_layout_idx = 1 if num_layouts > 1 else 0
        
    if financials_layout_idx >= num_layouts:
        print(f"Warning: Layout {financials_layout_idx} not found. Fallback to 2.")
        financials_layout_idx = 2 if num_layouts > 2 else 0       
    # Unified Text Selection: prefer business_overview, fallback to brand_overview
    overview_text = slides_data.get("business_overview") or slides_data.get("brand_overview", "")

    # Slide 1: Title Slide
    add_slide_from_template(prs, 0)
    
    # Slide 2: Company Profile
    add_slide_from_template(
        prs,
        layout_index=profile_layout_idx,
        placeholder_text={
           10: overview_text,
           11: slides_data.get("portfolio_and_products", ""),
           15: slides_data.get("at_a_glance", "")
        },
        placeholder_images={
            # Certifications
            12: cert[0], 
            13: cert[1], 
            14: cert[2],
            
            16: icon_list[0],  # Icon for business_overview
            17: icon_list[1],  # Icon for portfolio_and_products
            18: icon_list[2],  # Icon for at_a_glance
        }
    )
    
    # Slide 3: Financial Charts
    add_slide_from_template(
        prs,
        layout_index=financials_layout_idx,
        placeholder_text={
            12: slides_data.get("bar_chart_text", ""),
            13: slides_data.get("pie_chart_text", ""),
            18: f"**{slides_data.get('bar_chart_title', '')}**", 
            19: f"**{slides_data.get('pie_chart_title', '')}**",
            20: slides_data.get("financials_slide_title", ""),
        },
        placeholder_charts={
           10: (XL_CHART_TYPE.COLUMN_CLUSTERED, bar_data),
           11: (XL_CHART_TYPE.PIE, pie_data)
        },
        placeholder_images={

             17: icon_list[3],  # Icon for bar_chart
             16: icon_list[4],  # Icon for pie_chart
        }
    )
    
    # Slide 4: Investment Highlights
    highlights = slides_data.get("investment_highlights", [])
    # Pad with empty strings if fewer than 6
    while len(highlights) < 6:
        highlights.append("")

    add_slide_from_template(
        prs,
        layout_index=3,
        placeholder_text={
            10: highlights[0],
            11: highlights[1],
            12: highlights[2],
            13: highlights[3],
            14: highlights[4],
            15: highlights[5]
        },
        placeholder_images={
            16: os.path.join(IMAGE_DIR, slides_data.get("highlight_images", "Industry.jpeg").split("||")[0].strip())
        }
    )
    # Debug
    # print("Investment Highlight Image paths:")
    # for idx, img_path in enumerate(highlight_images):
    #     # print(f"  Raw {idx}: '{img_path}'")  # Show what we got
    #     if img_path:
    #         full_path = os.path.join(IMAGE_DIR, img_path.strip())
    #         exists = os.path.exists(full_path)
    #         # print(f"  Image {idx}: {full_path} - Exists: {exists}")

    # Slide 5: Disclaimer
    add_slide_from_template(prs, 4)
    
    # Delete the default slide master slide
    delete_slide(prs, 0)

    #Saves PPT
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_dir = os.path.join(project_root, "output", "ppt")
    os.makedirs(output_dir, exist_ok=True)
    output_filename = os.path.join(output_dir, f"company_{timestamp}.pptx")
    prs.save(output_filename)
